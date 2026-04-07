import io
import json
from pathlib import Path

import base64
import pandas as pd
import streamlit as st
import plotly.express as px
from scripts.plot_templates import kanban_counts, scorecard_kpis, sparkline_series, treemap_hierarchical, heatmap_matrix, stacked_bars, bubble_chart, sankey_chart, waterfall_chart, radar_chart, funnel_chart
import numpy as np
import plotly.graph_objects as go

try:
    import pptx
    from pptx.util import Inches
    PPTX_AVAILABLE = True
except Exception:
    PPTX_AVAILABLE = False

# ROOT del proyecto (padre de /scripts)
ROOT = Path(__file__).resolve().parents[1]

# Ruta archivo de consolidados
CONSOLIDADO_XLSX = ROOT / "data" / "output" / "Resultados Consolidados.xlsx"

@st.cache_data
def load_consolidado(path: Path) -> pd.DataFrame:
    # Preferir CSV ya normalizado si existe
    ts_csv = ROOT / 'data' / 'output' / 'analytics' / 'resultados_consolidados_timeseries.csv'
    latest_csv = ROOT / 'data' / 'output' / 'artifacts' / 'resultados_consolidados_latest.csv'
    if ts_csv.exists():
        df = pd.read_csv(ts_csv)
        # parse Fecha si existe
        if 'Fecha' in df.columns:
            df['Fecha'] = pd.to_datetime(df['Fecha'], errors='coerce')
        return df
    # fallback: read Excel path
    xl = pd.ExcelFile(path)
    sheet = xl.sheet_names[0]
    df = xl.parse(sheet)
    return df


@st.cache_data
def load_consolidado_latest() -> pd.DataFrame:
    latest_csv = ROOT / 'data' / 'output' / 'artifacts' / 'resultados_consolidados_latest.csv'
    xlsx = ROOT / 'data' / 'output' / 'Resultados Consolidados.xlsx'
    if latest_csv.exists():
        df = pd.read_csv(latest_csv)
        if 'Fecha' in df.columns:
            df['Fecha'] = pd.to_datetime(df['Fecha'], errors='coerce')
        return df
    if xlsx.exists():
        xl = pd.ExcelFile(xlsx)
        sheet = xl.sheet_names[0]
        df = xl.parse(sheet)
        if 'Fecha' in df.columns:
            df['Fecha'] = pd.to_datetime(df['Fecha'], errors='coerce')
        # take last per Id
        if 'Fecha' in df.columns:
            idx = df.groupby('Id')['Fecha'].idxmax()
            latest = df.loc[idx][['Id', 'Cumplimiento']]
            return latest
    return pd.DataFrame()

# Prototipo Nivel 3 — Kanban operativo y lista OM
# Supuestos: se usa `data/output/artifacts/indicadores_cmi_mapping_v2.csv`
# y `data/output/artifacts/ingesta_20260406_215409.json` como datos de apoyo.

ROOT = Path(__file__).resolve().parents[1]
MAPPING_CSV = ROOT / "data" / "output" / "artifacts" / "indicadores_cmi_mapping_v2.csv"
INGESTA_JSON = ROOT / "data" / "output" / "artifacts" / "ingesta_20260406_215409.json"


@st.cache_data
def load_mapping(path: Path) -> pd.DataFrame:
    df = pd.read_csv(path)
    return df


@st.cache_data
def load_ingesta(path: Path) -> dict:
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)


def derive_status(df: pd.DataFrame, green_thresh: float = 0.95, yellow_thresh: float = 0.8) -> pd.DataFrame:
    # Estado real usando cumplimiento del periodo más reciente
    # Requiere que df tenga columnas: Id, Cumplimiento
    def status_row(r):
        c = r.get("Cumplimiento")
        try:
            if pd.notna(c):
                c = float(c)
                if c >= green_thresh:
                    return "Actualizado"
                elif c >= yellow_thresh:
                    return "Pendiente"
                else:
                    return "Alerta"
        except Exception:
            pass
        return "Pendiente"

    df = df.copy()
    df["Status"] = df.apply(status_row, axis=1)
    return df


def main():
    st.set_page_config(layout="wide", page_title="Prototipo Nivel 3 — Operativo")
    st.title("Prototipo Nivel 3 — Operativo y Calidad")

    # Load data
    map_df = load_mapping(MAPPING_CSV)
    ingesta = load_ingesta(INGESTA_JSON)
    try:
        consol = load_consolidado(CONSOLIDADO_XLSX)
        # normalizar Fecha y Cumplimiento
        if isinstance(consol, pd.DataFrame) and not consol.empty:
            if 'Fecha' in consol.columns:
                consol['Fecha'] = pd.to_datetime(consol['Fecha'], errors='coerce')
            if 'Cumplimiento' in consol.columns:
                consol['Cumplimiento'] = pd.to_numeric(consol['Cumplimiento'], errors='coerce')
    except Exception:
        consol = None

    # Sidebar
    st.sidebar.header("Filtros")
    niveles = sorted(map_df["Nivel"].dropna().unique().tolist())
    nivel_sel = st.sidebar.selectbox("Nivel", options=niveles, index=max(0, niveles.index("Nivel 3") if "Nivel 3" in niveles else 0))
    periodicidades = sorted(map_df["Periodicidad"].dropna().unique().tolist())
    per_sel = st.sidebar.multiselect("Periodicidad", options=periodicidades, default=periodicidades)
    st.sidebar.markdown("---")
    st.sidebar.header("Parámetros y umbrales")
    green_thresh = st.sidebar.number_input("Umbral verde (>=)", min_value=0.0, max_value=1.0, value=0.95, step=0.01, format="%.2f")
    yellow_thresh = st.sidebar.number_input("Umbral amarillo (>=)", min_value=0.0, max_value=1.0, value=0.80, step=0.01, format="%.2f")
    rolling_window = st.sidebar.slider("Ventana rolling (periodos) para bandas", min_value=1, max_value=12, value=3)
    conf_mult = st.sidebar.number_input("Multiplicador banda confianza", value=1.28, min_value=0.1, max_value=10.0, step=0.01, format="%.2f")
    # Filtros adicionales solicitados: Unidad (Linea) y Proceso
    st.sidebar.markdown("---")
    st.sidebar.header("Filtros globales")
    unidades = sorted(map_df['Linea'].dropna().unique().astype(str).tolist()) if 'Linea' in map_df.columns else []
    unidad_sel = st.sidebar.multiselect('Unidad (Linea)', options=unidades, default=unidades)
    procesos = []
    if 'consol' in locals() and consol is not None and 'Proceso' in consol.columns:
        procesos = sorted(consol['Proceso'].dropna().unique().astype(str).tolist())
    proceso_sel = st.sidebar.multiselect('Proceso', options=procesos, default=procesos)

    # Inicializar df a partir del mapping para evitar UnboundLocalError
    df = map_df.copy()

    # Summary metrics
    col1, col2, col3 = st.columns(3)
    resumen = ingesta.get("resumen", {})
    col1.metric("Archivos procesados", resumen.get("total_archivos", "-"))
    col2.metric("Registros totales", resumen.get("total_registros", "-"))
    col3.metric("Archivos exitosos", resumen.get("exitosos", "-"))

    # Scorecard KPIs (dummy: % actualizados, % alerta, TMO OM)
    # Estos se recalculan más abajo después de aplicar filtros y merges con consol
    tmo_om = 3.5  # valor simulado
    # Mostrar placeholder de KPIs; se actualizan luego si procede
    st.plotly_chart(scorecard_kpis(
        pd.DataFrame({"Actualizados %": [0], "Alerta %": [0], "TMO OM": [tmo_om]}),
        ["Actualizados %", "Alerta %", "TMO OM"]
    ), use_container_width=True)

    # Filter indicators
    df = map_df[map_df["Nivel"] == nivel_sel]
    if per_sel:
        df = df[df["Periodicidad"].isin(per_sel)]
    # apply unidad filter
    if unidad_sel:
        df = df[df['Linea'].astype(str).isin(unidad_sel)]

    # Enlazar con cumplimiento real (por Id, periodo más reciente)
    if consol is not None and "Id" in consol.columns and "Cumplimiento" in consol.columns and "Fecha" in consol.columns:
        # Para cada Id, tomar el registro con la fecha más reciente
        # aplicar filtro de proceso a consol si aplica
        consol_filtered = consol
        if proceso_sel:
            consol_filtered = consol[consol['Proceso'].astype(str).isin(proceso_sel)]
        idx = consol_filtered.groupby("Id")['Fecha'].idxmax()
        consol_latest = consol_filtered.loc[idx][["Id","Cumplimiento"]]
        df = df.merge(consol_latest, on="Id", how="left")
    else:
        df["Cumplimiento"] = None

    df = derive_status(df, green_thresh=green_thresh, yellow_thresh=yellow_thresh)

    # KPI adicionales
    if consol is not None and 'Cumplimiento' in consol.columns:
        # usar consol_latest si disponible
        try:
            mean_cump = pd.to_numeric(consol_latest['Cumplimiento'], errors='coerce').dropna().mean()
            median_cump = pd.to_numeric(consol_latest['Cumplimiento'], errors='coerce').dropna().median()
            indicadores_alerta_pct = 100.0 * (df['Status'] == 'Alerta').sum() / len(df) if len(df) else 0
        except Exception:
            mean_cump = None
            median_cump = None
            indicadores_alerta_pct = 0
    else:
        mean_cump = None
        median_cump = None
        indicadores_alerta_pct = 0

    # Mostrar KPIs adicionales
    st.subheader('KPIs adicionales')
    k1, k2, k3 = st.columns(3)
    k1.metric('Cumplimiento global (media %)', f"{(mean_cump*100):.1f}%" if mean_cump is not None else '-')
    k2.metric('Cumplimiento mediano (media %)', f"{(median_cump*100):.1f}%" if median_cump is not None else '-')
    k3.metric('Indicadores en alerta (%)', f"{indicadores_alerta_pct:.1f}%")

    # Tendencia global de cumplimiento (serie agregada por fecha)
    if consol is not None and 'Fecha' in consol.columns and 'Cumplimiento' in consol.columns:
        ts = consol.copy()
        ts['Fecha'] = pd.to_datetime(ts['Fecha'], errors='coerce')
        ts = ts.dropna(subset=['Fecha'])
        # agregación mensual (media)
        ts_agg = ts.set_index('Fecha').resample('M')['Cumplimiento'].mean().dropna()
        if not ts_agg.empty:
            # banda de confianza simple: rolling mean +/- 1.28*rolling std (aprox 80%)
            roll_mean = ts_agg.rolling(rolling_window, min_periods=1).mean()
            roll_std = ts_agg.rolling(rolling_window, min_periods=1).std().fillna(0)
            upper = roll_mean + conf_mult * roll_std
            lower = roll_mean - conf_mult * roll_std
            fig = go.Figure()
            fig.add_trace(go.Scatter(x=ts_agg.index, y=ts_agg.values, mode='lines+markers', name='Cumplimiento'))
            fig.add_trace(go.Scatter(x=roll_mean.index, y=roll_mean.values, mode='lines', name='Rolling mean'))
            fig.add_trace(go.Scatter(x=upper.index, y=upper.values, fill=None, mode='lines', line=dict(color='rgba(255,0,0,0.1)'), showlegend=False))
            fig.add_trace(go.Scatter(x=lower.index, y=lower.values, fill='tonexty', mode='lines', line=dict(color='rgba(255,0,0,0.1)'), name='80% band'))
            fig.update_layout(title='Tendencia agregada de cumplimiento', yaxis_title='Cumplimiento (ratio)')
            st.plotly_chart(fig, use_container_width=True)

    # Selector para visualizar histórico de un indicador
    st.sidebar.markdown("---")
    indicador_options = df[['Id', 'Indicador']].drop_duplicates().head(200)
    indicador_map = {f"{r.Id} - {r.Indicador}": r.Id for _, r in indicador_options.iterrows()}
    sel_key = st.sidebar.selectbox("Mostrar histórico de indicador", options=[""] + list(indicador_map.keys()))
    sel_id = indicador_map.get(sel_key)

    if sel_id and consol is not None:
        hist = consol[consol['Id'] == sel_id].sort_values('Fecha')
        if not hist.empty and 'Cumplimiento' in hist.columns:
            fig_sp = sparkline_series(hist['Cumplimiento'])
            st.subheader(f"Histórico — Indicador {sel_id}")
            st.plotly_chart(fig_sp, use_container_width=True)
        else:
            st.info('No hay serie histórica disponible para el indicador seleccionado.')

    # Kanban gráfico
    st.subheader("Kanban — Estado de indicadores (prototipo)")
    st.plotly_chart(kanban_counts(df), use_container_width=True)
    kcols = st.columns(3)
    for i, state in enumerate(["Actualizado", "Pendiente", "Alerta"]):
        with kcols[i]:
            d_state = df[df["Status"] == state] if state != "Alerta" else df.sample(frac=0.05) if not df.empty else df
            st.markdown(f"**{state}** — {len(d_state)}")
            for _, r in d_state.head(10).iterrows():
                st.write(f"- {r.get('Id')} — {r.get('Indicador')}")

    st.subheader("Detalle de indicadores")
    st.dataframe(df[["Id", "Indicador", "Linea", "Periodicidad", "Meta", "Orden CMI", "Status"]].reset_index(drop=True))

    # Expanders por fila con sparkline histórico (muestra hasta 20 filas)
    st.subheader("Detalle expandible por indicador")
    for _, r in df.head(20).iterrows():
        id_val = r.get('Id')
        title = f"{id_val} - {r.get('Indicador')}"
        with st.expander(title, expanded=False):
            if 'consol' in locals() and consol is not None:
                # comparar como strings para evitar desajustes de tipos
                hist = consol[consol['Id'].astype(str) == str(id_val)].sort_values('Fecha')
                if not hist.empty and 'Cumplimiento' in hist.columns:
                    st.plotly_chart(sparkline_series(hist['Cumplimiento']), use_container_width=True)
                    st.write(f"Último cumplimiento: {hist['Cumplimiento'].iloc[-1]}")
                else:
                    st.info('No hay serie histórica disponible para este indicador.')
            else:
                st.info('Archivo de consolidados no disponible.')

    # Export controls
    st.markdown("---")
    st.subheader("Exportar datos filtrados")
    csv_bytes = df.to_csv(index=False).encode('utf-8')
    st.download_button("Descargar CSV (filtrado)", data=csv_bytes, file_name="indicadores_filtrados.csv", mime="text/csv")

    # PPTX export (limitado)
    if PPTX_AVAILABLE:
        max_slides = st.number_input("Slides en PPTX (máx)", min_value=1, max_value=50, value=10)
        if st.button("Exportar PPTX (filtrado)"):
            prs = pptx.Presentation()
            blank_slide_layout = prs.slide_layouts[6]
            added = 0
            for _, row in df.head(max_slides).iterrows():
                slide = prs.slides.add_slide(blank_slide_layout)
                title = slide.shapes.title if slide.shapes.title else None
                # Add title as textbox
                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
                tf = txBox.text_frame
                tf.text = f"{row.get('Id')} - {row.get('Indicador')}"
                # generate sparkline image
                if 'consol' in locals() and consol is not None:
                    hist = consol[consol['Id'].astype(str) == str(row.get('Id'))].sort_values('Fecha')
                else:
                    hist = None
                try:
                    if hist is not None and not hist.empty:
                        fig = sparkline_series(hist['Cumplimiento'])
                        img_bytes = fig.to_image(format='png')
                        img_stream = io.BytesIO(img_bytes)
                        slide.shapes.add_picture(img_stream, Inches(1), Inches(1), width=Inches(8))
                except Exception:
                    pass
                added += 1
            out = io.BytesIO()
            prs.save(out)
            out.seek(0)
            st.download_button("Descargar PPTX", data=out, file_name="indicadores_filtrados.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
    else:
        st.info('Para exportar a PPTX instala `python-pptx` en el entorno.')

    # Inline sparklines in table view
    st.markdown("---")
    inline = st.checkbox('Mostrar sparklines inline en tabla (top 50)')
    if inline:
        st.subheader('Tabla con sparklines (muestra limitada)')
        rows_to_show = min(50, len(df))
        for _, r in df.head(rows_to_show).iterrows():
            c1, c2, c3, c4 = st.columns([1, 5, 1, 3])
            c1.write(r.get('Id'))
            c2.write(r.get('Indicador'))
            c3.write(r.get('Status'))
            with c4:
                if 'consol' in locals() and consol is not None:
                    hist = consol[consol['Id'].astype(str) == str(r.get('Id'))].sort_values('Fecha')
                    if not hist.empty and 'Cumplimiento' in hist.columns:
                        st.plotly_chart(sparkline_series(hist['Cumplimiento']), use_container_width=True)
                    else:
                        st.write('-')
                else:
                    st.write('-')

    st.subheader("Ingesta — detalle de últimos procesos")
    for d in ingesta.get("detalle", [])[:10]:
        st.write(f"**{d.get('archivo')}** — {d.get('plantilla')} — registros leídos: {d.get('registros_leidos')} — exitosa: {d.get('exitosa')}")

    # Analítica avanzada: selector de plantillas gráficas
    st.markdown("---")
    st.subheader('Analítica avanzada')
    chart_type = st.selectbox('Tipo de visualización', options=['Treemap','Heatmap','Barras apiladas','Bubble','Sankey','Waterfall','Radar','Funnel'], index=0)

    if st.button('Generar visualización'):
        try:
            if chart_type == 'Treemap':
                fig = treemap_hierarchical(map_df, path=['Linea','Clasificacion','Indicador'], values_col=None)
                st.plotly_chart(fig, use_container_width=True)
            elif chart_type == 'Heatmap':
                if consol is not None:
                    # pivot: Proceso x Mes -> Cumplimiento mean
                    if 'Mes' in consol.columns and 'Proceso' in consol.columns and 'Cumplimiento' in consol.columns:
                        fig = heatmap_matrix(consol, index='Proceso', columns='Mes', values='Cumplimiento')
                        st.plotly_chart(fig, use_container_width=True)
                    else:
                        st.info('Archivo consolidados no tiene columnas Proceso/Mes/Cumplimiento')
                else:
                    st.info('Consolidado no disponible')
            elif chart_type == 'Barras apiladas':
                fig = stacked_bars(map_df, x='Linea', y='Orden CMI', color='Clasificacion')
                st.plotly_chart(fig, use_container_width=True)
            elif chart_type == 'Bubble':
                # bubble: count indicadores por Linea
                counts = map_df.groupby('Linea').size().reset_index(name='count')
                fig = bubble_chart(counts, x='Linea', y='count', size='count')
                st.plotly_chart(fig, use_container_width=True)
            elif chart_type == 'Sankey':
                # Sankey using Linea -> Clasificacion
                sank_df = map_df[['Linea','Clasificacion']].copy()
                sank_df['value'] = 1
                sank_df = sank_df.groupby(['Linea','Clasificacion']).sum().reset_index()
                fig = sankey_chart(sank_df, source_col='Linea', target_col='Clasificacion', value_col='value')
                st.plotly_chart(fig, use_container_width=True)
            elif chart_type == 'Waterfall':
                # sample waterfall from top 10 Linea contributions
                agg = map_df.groupby('Linea').size().reset_index(name='count').sort_values('count', ascending=False).head(10)
                fig = waterfall_chart(agg, x='Linea', y='count')
                st.plotly_chart(fig, use_container_width=True)
            elif chart_type == 'Radar':
                # radar from average cumplimiento per Linea (if available)
                if consol is not None and 'Cumplimiento' in consol.columns and 'Proceso' in consol.columns:
                    agg = consol.groupby('Proceso')['Cumplimiento'].mean().dropna().reset_index()
                    cats = agg['Proceso'].astype(str).tolist()[:10]
                    vals = agg['Cumplimiento'].tolist()[:10]
                    fig = radar_chart(pd.DataFrame(), categories=cats, values=vals)
                    st.plotly_chart(fig, use_container_width=True)
                else:
                    st.info('No hay datos para radar')
            elif chart_type == 'Funnel':
                # dummy funnel from top processes
                agg = map_df.groupby('Linea').size().reset_index(name='count').sort_values('count', ascending=False).head(6)
                fig = funnel_chart(agg, x='count', y='Linea')
                st.plotly_chart(fig, use_container_width=True)
        except Exception as e:
            st.error(f'Error generando visualización: {e}')


if __name__ == "__main__":
    main()
